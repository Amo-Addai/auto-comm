from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Create a new presentation
prs = Presentation()

# Define custom colors (approximated from common pitch deck designs)
PRIMARY_COLOR = RGBColor(0, 51, 102)  # Dark blue
SECONDARY_COLOR = RGBColor(0, 153, 204)  # Light blue
ACCENT_COLOR = RGBColor(255, 102, 0)  # Orange
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray
LIGHT_TEXT = RGBColor(255, 255, 255)  # White

# Define fonts
TITLE_FONT = "Arial Bold"
SUBTITLE_FONT = "Arial"
BODY_FONT = "Arial"

# Function to add a title slide
def add_title_slide(title, subtitle=None):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(44)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.add_paragraph()
        p.text = subtitle
        p.font.name = SUBTITLE_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
    
    return slide

# Function to add a section header slide
def add_section_header_slide(title):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(54)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# Function to add a content slide with title and bullet points
def add_content_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Set content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear existing content
    
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(12)
    
    return slide

# Function to add a two-column slide
def add_two_column_slide(title, left_content, right_content):
    slide_layout = prs.slide_layouts[3]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Left column
    left_shape = slide.placeholders[1]
    left_tf = left_shape.text_frame
    left_tf.clear()
    
    for item in left_content:
        p = left_tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(8)
    
    # Right column
    right_shape = slide.placeholders[2]
    right_tf = right_shape.text_frame
    right_tf.clear()
    
    for item in right_content:
        p = right_tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(8)
    
    return slide

# Function to add a three-column slide
def add_three_column_slide(title, col1_content, col2_content, col3_content):
    slide_layout = prs.slide_layouts[2]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Define column positions
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(3)
    height = Inches(5)
    
    # Column 1
    col1_box = slide.shapes.add_textbox(left, top, width, height)
    col1_tf = col1_box.text_frame
    col1_tf.word_wrap = True
    
    for item in col1_content:
        p = col1_tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(8)
    
    # Column 2
    col2_box = slide.shapes.add_textbox(left + width + Inches(0.5), top, width, height)
    col2_tf = col2_box.text_frame
    col2_tf.word_wrap = True
    
    for item in col2_content:
        p = col2_tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(8)
    
    # Column 3
    col3_box = slide.shapes.add_textbox(left + 2*(width + Inches(0.5)), top, width, height)
    col3_tf = col3_box.text_frame
    col3_tf.word_wrap = True
    
    for item in col3_content:
        p = col3_tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(8)
    
    return slide

# Function to add a table slide
def add_table_slide(title, table_data):
    slide_layout = prs.slide_layouts[5]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Define table position and size
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4)
    
    # Add table
    table = slide.shapes.add_table(rows=len(table_data), cols=len(table_data[0]), 
                                  left=left, top=top, width=width, height=height).table
    
    # Set table content
    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.name = BODY_FONT
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
            
            # Header row styling
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
                cell.text_frame.paragraphs[0].font.color.rgb = LIGHT_TEXT
                cell.text_frame.paragraphs[0].font.bold = True
    
    return slide

# Function to add a timeline slide
def add_timeline_slide(title, timeline_data):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Add timeline
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1)
    
    for i, (period, content) in enumerate(timeline_data):
        # Period
        period_box = slide.shapes.add_textbox(left + i * (width + Inches(0.5)), top, width, height)
        period_frame = period_box.text_frame
        period_frame.word_wrap = True
        p = period_frame.add_paragraph()
        p.text = period
        p.font.name = TITLE_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Content
        content_box = slide.shapes.add_textbox(left + i * (width + Inches(0.5)), top + height, width, Inches(3))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.add_paragraph()
        p.text = content
        p.font.name = BODY_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# Function to add a financial summary slide
def add_financial_slide(title, financial_data, deal_data):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Financial summary
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(4)
    height = Inches(3)
    
    financial_box = slide.shapes.add_textbox(left, top, width, height)
    financial_frame = financial_box.text_frame
    financial_frame.word_wrap = True
    
    for item in financial_data:
        p = financial_frame.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)
    
    # Deal information
    deal_box = slide.shapes.add_textbox(left + width + Inches(1), top, width, height)
    deal_frame = deal_box.text_frame
    deal_frame.word_wrap = True
    
    for item in deal_data:
        p = deal_frame.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)
    
    return slide

# Function to add a contact slide
def add_contact_slide(name, title, email, website, address):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add name and title
    name_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    name_frame = name_box.text_frame
    name_frame.word_wrap = True
    p = name_frame.add_paragraph()
    p.text = name
    p.font.name = TITLE_FONT
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = title
    p.font.name = SUBTITLE_FONT
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Add contact information
    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(4)
    height = Inches(0.5)
    
    # Email
    email_box = slide.shapes.add_textbox(left, top, width, height)
    email_frame = email_box.text_frame
    email_frame.word_wrap = True
    p = email_frame.add_paragraph()
    p.text = f"✉ {email}"
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Website
    website_box = slide.shapes.add_textbox(left, top + height, width, height)
    website_frame = website_box.text_frame
    website_frame.word_wrap = True
    p = website_frame.add_paragraph()
    p.text = f" Visit Us {website}"
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Address
    address_box = slide.shapes.add_textbox(left, top + 2*height, width, height)
    address_frame = address_box.text_frame
    address_frame.word_wrap = True
    p = address_frame.add_paragraph()
    p.text = f" {address}"
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Add "Questions?" text
    questions_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    questions_frame = questions_box.text_frame
    questions_frame.word_wrap = True
    p = questions_frame.add_paragraph()
    p.text = "Questions?"
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add icons
    icons = ["Investors", "Partners", "Projects", "Advisors", "Markets", "Impact"]
    left = Inches(1.5)
    top = Inches(7)
    width = Inches(1.5)
    height = Inches(0.5)
    
    for i, icon in enumerate(icons):
        icon_box = slide.shapes.add_textbox(left + i * (width + Inches(0.5)), top, width, height)
        icon_frame = icon_box.text_frame
        icon_frame.word_wrap = True
        p = icon_frame.add_paragraph()
        p.text = icon
        p.font.name = BODY_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Add "Schedule a Demo" button
    demo_box = slide.shapes.add_textbox(Inches(3.5), Inches(8), Inches(3), Inches(0.7))
    demo_frame = demo_box.text_frame
    demo_frame.word_wrap = True
    p = demo_frame.add_paragraph()
    p.text = "Schedule a Demo "
    p.font.name = TITLE_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Format button background
    fill = demo_box.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR
    
    return slide

# Slide 1: Title Slide
slide1 = add_title_slide(
    "Institutional Investment Platform",
    "Digital-first platform empowering institutional investors to structure,\nmanage, and deploy capital into frontier-market projects through data-\ndriven co-investment infrastructure.\n\nOur Vision 〉"
)

# Slide 2: Platform Description
slide2 = add_content_slide(
    "CoFund is a digital-first platform that empowers institutional investors to structure, manage, and deploy capital into frontier-market projects through data-driven co-investment infrastructure.",
    [
        " Investors",
        " Fund Managers",
        " Projects",
        " Infrastructure",
        " Compliance",
        " Data",
        " CoFund"
    ]
)

# Slide 3: The Pain
slide3 = add_three_column_slide(
    " The Pain",
    [
        " Fragmented Deal Flow",
        "Institutional capital struggles to find quality opportunities",
        "",
        " Manual Fund Structuring",
        "Time-consuming processes with high error potential",
        "",
        " Disconnected Tools",
        "ESG, risk, and compliance systems operate in silos",
        "",
        " Minimal Transparency",
        "Limited visibility post-deployment"
    ],
    [
        " Today's Workaround",
        "",
        " Excel Models",
        "Prone to errors and difficult to maintain",
        "",
        " Legal PDFs",
        "Static documents with limited functionality",
        "",
        " Disconnected Data Rooms",
        "Fragmented information sources",
        "",
        "✉ Siloed Communications",
        "Inefficient investor updates and reporting"
    ],
    [
        " The Result?",
        "",
        "Lost Deals",
        "Lower Trust",
        "Inefficiency",
        "Missed Returns",
        "",
        "          ⚠ "
    ]
)

# Slide 4: CoFund Solution
slide4 = add_two_column_slide(
    "CoFund provides a unified digital platform to manage every step of institutional project-based financing:",
    [
        " Fund Structuring Engine",
        "Create SPVs, LPs, syndicates in minutes",
        "",
        " Project Marketplace",
        "Curated, ESG-qualified opportunities",
        "",
        " Risk & ESG Engine",
        "Real-time impact + risk dashboards",
        "",
        " Returns Forecasting",
        "ML-based modeling on deal performance",
        "",
        " Capital Dashboard",
        "Investor dashboards, cap tables",
        "",
        " Compliance Module",
        "KYC/AML, SEC/FCA workflows"
    ],
    [
        " The Stripe for Frontier Market Investment",
        "Securely digitizing the entire co-investment stack",
        "",
        " Compare with Stripe",
        "",
        " Fund Structuring",
        " Project Marketplace",
        " Risk & ESG",
        " Returns",
        " Capital Monitor",
        " Compliance",
        "",
        " CoFund"
    ]
)

# Slide 5: Timing Shift
slide5 = add_content_slide(
    " Timing Shift",
    [
        "Emerging markets now represent the largest untapped investment frontier",
        "",
        " Global Trends",
        "ESG mandates, infrastructure development, green finance",
        "",
        " Tech Evolution",
        "APIs, digital onboarding, embedded compliance",
        "",
        " Market Frustration",
        "Investors demand faster, smarter, more transparent vehicles",
        "",
        " The time to digitize is now",
        "CoFund is purpose-built for this moment",
        "",
        "The convergence of market trends creates the perfect moment for CoFund",
        "",
        "Past 2010s 2020s Now Future"
    ]
)

# Slide 6: Market Focus
slide6 = add_content_slide(
    " Targeting $15T+ Global Market Frontier Markets Focus",
    [
        " Who We Serve",
        "",
        " Institutional Investors",
        "PE, VC, sovereign funds, family offices",
        "",
        " Fund Managers & Syndicates",
        "",
        " Banks & Mutual Funds",
        "expanding Africa exposure",
        "",
        " Infrastructure Funds",
        "$1.3T+ market",
        "",
        " PE Co-investment platforms",
        "",
        " ESG Green Funds",
        "compliant investments",
        "",
        "TAM SAM SOM",
        "",
        " Global Private Capital AUM",
        "$18T",
        "Growing market with increasing interest in frontier opportunities",
        "",
        " $1B AUM Goal (Yr 3)",
        " Just 0.01% of global TAM, unlocking massive economic transformation"
    ]
)

# Slide 7: Product Roadmap
slide7 = add_two_column_slide(
    " Product Roadmap",
    [
        " Fund Builder",
        "SPVs, LPs, JV structures",
        "",
        " Deal Marketplace",
        "Curated, verified projects",
        "",
        " ESG/Risk Engine",
        "Real-time analytics",
        "",
        " Forecasting Tool",
        "ML-based return modeling",
        "",
        " Capital Dashboard",
        "Investor dashboards, cap tables",
        "",
        " Compliance Module",
        "KYC/AML, SEC/FCA workflows"
    ],
    [
        " Technology Architecture",
        "",
        "☁ SaaS, Cloud-Native",
        " API-First",
        " Optional Blockchain Ledger",
        " GDPR/BoG/SEC-Compliant",
        "",
        " Product Roadmap",
        "",
        "Q1: MVP Launch",
        "5 Anchor Clients",
        "",
        "Q2: ESG AI Engine",
        "Beta + Regulatory Integrations",
        "",
        "Q3: Deal Flow Expansion",
        "50+ Projects",
        "",
        "Q4: Syndicate Tooling",
        "Co-investment Capabilities",
        "",
        " Q1"
    ]
)

# Slide 8: Revenue Streams
slide8 = add_two_column_slide(
    " Revenue Streams",
    [
        " Platform Subscription",
        "Tiered SaaS for GPs, LPs, fund admins",
        "",
        "% AUM Fees",
        "0.25%–1% of assets under management",
        "",
        " Data Tools",
        "Forecasting & ESG scoring modules (premium)",
        "",
        " Transaction Fees",
        "Flat per-deployment event",
        "",
        "© Licensing",
        "For fund admins or custodians (white-label)"
    ],
    [
        " LTV (Avg Account Lifetime Value)",
        "",
        " PE/VC Clients",
        "$120K–$400K/year",
        "",
        " Smaller Fund Managers",
        "$20K–$50K/year",
        "",
        " Licensing",
        "$1M+ ARR",
        "",
        " CAC (Avg Customer Acquisition Cost)",
        " [TODO]",
        "Customer acquisition metrics to be determined",
        "",
        "⚖ CAC:LTV Ratio",
        "",
        " Strong ratio gap indicates high profitability potential",
        " Recurring revenue model enhances customer lifetime value",
        " Platform approach increases stickiness and reduces churn"
    ]
)

# Slide 9: Competition
slide9 = add_table_slide(
    "Player Focus Weakness CoFund's Edge",
    [
        ["Player", "Focus", "Weakness", "CoFund's Edge"],
        ["Carta", "Equity/Cap Table", "No project finance workflows", "Africa-first"],
        ["Juniper Square", "Fund Admin", "Not ESG-native, lacks project modeling", "ESG-optimized"],
        ["iCapital", "Alt Access", "Closed ecosystem, not co-investment-ready", "Open platform"],
        ["AngelList", "Startup syndicates", "US/early-stage only, not institutional-grade", "Institutional-grade"]
    ]
)

# Add CoFund's Edge section
slide9 = add_content_slide(
    " CoFund's Edge",
    [
        " Africa-first",
        " ESG-compliant",
        " Project-finance optimized",
        " Regulatory-grade transparency",
        " Built for both institutions and emerging managers",
        " Modular pricing and onboarding",
        "",
        " Any Potential Moats?",
        "",
        " Data Network Effects",
        "Growing data advantage as more projects and investors join the platform",
        "",
        " Partnership Ecosystem",
        "Strategic alliances with key market players create barriers to entry",
        "",
        " Regulatory Expertise",
        "Deep compliance knowledge in emerging markets creates switching costs"
    ]
)

# Slide 10: Channels
slide10 = add_content_slide(
    " Channels",
    [
        " Direct Sales",
        "Asset managers & PE firms",
        "",
        " Ecosystem Referrals",
        "Mest, KIC, Impact Hub",
        "",
        " Co-marketing",
        "Fund admins & law firms",
        "",
        " Thought Leadership",
        "ESG capital, digital investing",
        "",
        " Launch Activities",
        "",
        " Target: $50M AUM in Year 1",
        "",
        "Private Pilot",
        " 5 PE/VC funds",
        "",
        "Sponsored Webinars",
        " ESG + digital investing",
        "",
        "Roadshow",
        " US-Africa investment summits",
        "",
        "LP Roundtables",
        " Diaspora & sovereign fund capital",
        "",
        "Strategic channel partnerships driving rapid adoption"
    ]
)

# Slide 11: Core Team
slide11 = add_content_slide(
    " Core Team",
    [
        " Kwadwo Amo-Addai",
        "CEO/Founder",
        "Fund Strategy Africa Project Finance",
        "",
        " [TBD]",
        "CTO",
        "Platform Architecture Security",
        "",
        " [TBD]",
        "CIO & Head of Investments",
        "Deal Pipeline Investor Network",
        "",
        " [TBD]",
        "Risk & ESG Lead",
        "ESG Modeling Compliance",
        "",
        " Advisory Board",
        "",
        " Ex-Bankers",
        "Institutional expertise",
        "",
        "⚖ Legal Experts",
        "Regulatory guidance",
        "",
        " Fintech Leaders",
        "Technology innovation",
        "",
        " Market Specialists",
        "Frontier market insights",
        "",
        " Supported By"
    ]
)

# Slide 12: Financials Summary
slide12 = add_financial_slide(
    " Financials Summary",
    [
        "Y1 Revenue $1.2M",
        "Y2 Revenue $6.8M",
        "Y3 Revenue $20M+",
        "Break-even 18 months",
        "Gross Margin 80%+ (Y2+)",
        "",
        "View Detailed Financial Model"
    ],
    [
        " The Deal",
        "",
        " Raise: $5M Pre-Seed",
        "Convertible note, $20M cap / 20% discount",
        "",
        "Use of Funds:",
        "Product Development 40%",
        "Market Growth 30%",
        "Legal & Compliance 15%",
        "ESG/Risk Analytics 15%",
        "",
        " Target close: [insert date]",
        "45 days remaining",
        "",
        " Vision & Impact",
        "",
        "Building digital rails for capital movement in emerging markets",
        "",
        " Unlock generational returns through smarter capital flow",
        "",
        " Move from extractive to collaborative capital models",
        "",
        " Bring regulatory-grade transparency to frontier investing",
        "",
        " Ghana →  Africa →  Global →  Impact"
    ]
)

# Slide 13: Contact
slide13 = add_contact_slide(
    "Kwadwo Amo-Addai",
    "CEO/Founder",
    "kwadwoamoad@gmail.com",
    "www.cofund.io",
    "Accra, Ghana"
)

# Save the presentation
prs.save('CoFund_Pitch_Deck.pptx')
print("Pitch deck created successfully!")
